#!/usr/bin/env python3
"""Structured PPTX-to-Markdown extractor.

Produces a Markdown representation of a .pptx deck that preserves spatial
layout, connector topology, group nesting, and shape-type semantics — the
structural information lost in flat text extractions (e.g. markitdown).

Usage:
    python3 extract_pptx_structured.py input.pptx
    python3 extract_pptx_structured.py input.pptx -o output.md
    python3 extract_pptx_structured.py input.pptx --config /path/to/config.yaml

Reads:  input.pptx (any PowerPoint file)
Writes: <input_stem>_structured.md (same dir, or -o path)

Config: reads config.yaml from the skill directory (or --config path).
        Falls back to built-in defaults if config is missing.

Dependencies: python-pptx, lxml (pip install python-pptx lxml)
"""

import argparse
import base64
import math
import os
import platform
import re
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

# ── Slide image export ─────────────────────────────────────────────────────

def export_slide_images(pptx_path, output_dir, slide_nums=None, width=1920, height=1080):
    """Export slides as PNG images using PowerPoint COM (Windows) or LibreOffice (Linux).

    Returns a dict {slide_num: relative_image_path} or {} if export fails.
    Images are saved as slide_NN.png in output_dir.
    """
    images = {}
    is_windows = platform.system() == "Windows"

    if is_windows:
        ppt_app = None
        pres = None
        try:
            import comtypes.client
            comtypes.CoInitialize()
            try:
                ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
                ppt_app.Visible = 1
                abs_path = os.path.abspath(pptx_path)
                pres = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)

                if slide_nums is None:
                    slide_nums = list(range(1, pres.Slides.Count + 1))

                for sn in slide_nums:
                    out_path = os.path.join(output_dir, f"slide_{sn:02d}.png")
                    pres.Slides(sn).Export(out_path, "PNG", width, height)
                    images[sn] = f"slide_{sn:02d}.png"
            finally:
                if pres is not None:
                    try:
                        pres.Close()
                    except Exception:
                        pass
                if ppt_app is not None:
                    try:
                        ppt_app.Quit()
                    except Exception:
                        pass
                comtypes.CoUninitialize()
        except Exception as e:
            print(f"  (COM export failed: {e}; trying LibreOffice fallback)", file=sys.stderr)
            images = _export_via_libreoffice(pptx_path, output_dir, slide_nums)
    else:
        images = _export_via_libreoffice(pptx_path, output_dir, slide_nums)

    return images


def _export_via_libreoffice(pptx_path, output_dir, slide_nums=None):
    """Fallback: use LibreOffice to convert to PDF, then pdftoppm to PNG."""
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
                check=True, capture_output=True, timeout=120
            )
            pdf_name = Path(pptx_path).stem + ".pdf"
            pdf_path = os.path.join(tmpdir, pdf_name)

            result = subprocess.run(
                ["pdftoppm", "-jpeg", "-r", "150", pdf_path,
                 os.path.join(output_dir, "slide")],
                capture_output=True
            )

            if result.returncode != 0:
                return {}

            images = {}
            for f in os.listdir(output_dir):
                if f.startswith("slide") and f.endswith((".png", ".jpg", ".jpeg")):
                    m = re.search(r'slide[-_]?(\d+)', f)
                    if m:
                        sn = int(m.group(1))
                        # Only return requested slides (or all if slide_nums is None)
                        if slide_nums is None or sn in slide_nums:
                            images[sn] = f
            return images
    except Exception as e:
        print(f"  (LibreOffice export also failed: {e})", file=sys.stderr)
        return {}


def embed_image_base64(image_path, max_size_kb=500):
    """Read an image file and return a base64 data URI, or None if too large."""
    try:
        file_size = os.path.getsize(image_path)
        if file_size > max_size_kb * 1024:
            return None
        with open(image_path, "rb") as f:
            data = base64.b64encode(f.read()).decode("ascii")
        ext = Path(image_path).suffix.lstrip(".")
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg"}.get(ext, "image/png")
        return f"data:{mime};base64,{data}"
    except Exception:
        return None


# ── Defaults (used when config.yaml is absent) ───────────────────────────

DEFAULTS = {
    "output_dir": "",
    "output_suffix": "_structured.md",
    "include_deck_summary": True,
    "include_notes": True,
    "truncate_text": 200,
    "embed_slide_images": True,
    "image_embed_threshold": 10,
    "image_embed_max_kb": 500,
    # Legacy keys (used by dead code kept for potential future use)
    "diagram_heavy_threshold": 15,
    "row_cluster_threshold": 0.45,
    "col_cluster_threshold": 0.45,
    "connector_proximity": 0.35,
    "grid_min_shapes": 4,
    "grid_min_rows": 2,
}

# ── Shape-type → semantic label mapping ──────────────────────────────────

PRST_GEO_TO_LABEL = {
    "rect": "box", "roundRect": "step", "round2SameRect": "step",
    "round2DiagRect": "step", "diamond": "decision", "hexagon": "phase",
    "chevron": "phase", "arrow": "arrow", "rightArrow": "arrow",
    "leftArrow": "arrow", "upArrow": "arrow", "downArrow": "arrow",
    "bentArrow": "arrow", "curvedArrow": "arrow", "pentagon": "phase",
    "parallelogram": "data", "trapezoid": "data", "oval": "node",
    "ellipse": "node", "cylinder": "store", "cloud": "cloud",
    "calloutRect": "note", "wedgeRectCallout": "note",
    "wedgeRoundRectCallout": "note", "wedgeEllipseCallout": "note",
    "cloudCallout": "note", "ribbon": "tag", "ribbon2": "tag",
    "flowChartProcess": "process", "flowChartDecision": "decision",
    "flowChartTerminator": "terminator", "flowChartConnector": "connector",
    "flowChartDocument": "document", "flowChartInputOutput": "io",
    "flowChartPredefinedProcess": "predefined", "flowChartInternalStorage": "store",
    "flowChartExtract": "extract", "flowChartMerge": "merge",
    "flowChartSort": "sort", "flowChartOr": "or",
    "flowChartSummingJunction": "junction", "flowChartManualInput": "manual-input",
    "flowChartManualOperation": "manual-op", "flowChartPreparation": "prep",
    "flowChartDisplay": "display", "flowChartStoredData": "stored-data",
    "flowChartSequentialAccessStorage": "seq-store",
    "flowChartMagneticDisk": "disk", "flowChartDirectAccessStorage": "direct-store",
    "flowChartOfflineStorage": "offline-store", "flowChartPunchedCard": "card",
    "flowChartPunchedTape": "tape", "flowChartCollate": "collate",
    "frame": "frame", "plaque": "marker", "star5": "marker",
    "star4": "marker", "star6": "marker", "star8": "marker",
    "plus": "marker", "mathPlus": "marker", "mathMinus": "marker",
    "mathMultiply": "marker", "mathDivide": "marker",
    "mathEqual": "marker", "mathNotEqual": "marker",
    "bracePair": "bracket", "bracketPair": "bracket",
    "stopSign": "stop", "noSmoking": "forbidden", "doNotEnter": "forbidden",
    "heart": "deco", "lightningBolt": "deco", "smileyFace": "deco",
    "sun": "deco", "moon": "deco", "bevel": "deco",
    "homePlate": "phase", "decagon": "node", "heptagon": "node",
    "octagon": "node", "nonagon": "node", "pie": "node",
    "pieWedge": "node", "teardrop": "node",
    "circularArrow": "arrow", "circularText": "text",
    "uturnArrow": "arrow", "leftUpArrow": "arrow",
    "leftRightUpArrow": "arrow", "quadArrow": "arrow",
    "leftRightArrow": "arrow", "upDownArrow": "arrow",
    "leftRightCircularArrow": "arrow", "swooshArrow": "arrow",
    "stripedRightArrow": "arrow", "notchedRightArrow": "arrow",
    "blockArc": "arc", "pie1": "arc", "pie2": "arc", "arc": "arc",
    "chord": "arc", "epoch": "tag",
    "borderCallout1": "note", "borderCallout2": "note", "borderCallout3": "note",
    "accentCallout1": "note", "accentCallout2": "note", "accentCallout3": "note",
    "explosion1": "marker", "explosion2": "marker",
    "leftBrace": "bracket", "leftBracket": "bracket",
    "rightBrace": "bracket", "rightBracket": "bracket",
    "brace": "bracket", "bracket": "bracket",
}


# ── Config loading ───────────────────────────────────────────────────────

def load_config(config_path: Optional[str] = None) -> dict:
    """Load config from YAML, falling back to defaults."""
    cfg = dict(DEFAULTS)

    # Try to find config.yaml relative to this script (skill dir)
    if config_path is None:
        script_dir = Path(__file__).parent.parent
        default_config = script_dir / "config.yaml"
        if default_config.exists():
            config_path = str(default_config)

    if config_path and os.path.exists(config_path):
        with open(config_path, "r", encoding="utf-8") as f:
            user_cfg = yaml.safe_load(f)
        if user_cfg:
            cfg.update(user_cfg)

    return cfg


# ── Data classes ─────────────────────────────────────────────────────────

@dataclass
class ShapeInfo:
    idx: int
    shape_type: str
    prst_geo: Optional[str]
    text: str
    left: float
    top: float
    width: float
    height: float
    cx: float
    cy: float
    group_path: str
    has_arrowhead: bool
    arrow_dir: str
    line_start: Optional[tuple]
    line_end: Optional[tuple]
    picture_name: str
    is_table: bool
    table_rows: list
    raw_shape: object = field(repr=False, default=None)


# ── Helpers ──────────────────────────────────────────────────────────────

def emu_to_inches(emu):
    if emu is None:
        return 0.0
    return Emu(emu).inches


def get_shape_text(sh):
    """Extract all text from a shape's text_frame, preserving line breaks."""
    if hasattr(sh, "has_text_frame") and not sh.has_text_frame:
        return ""
    if not hasattr(sh, "text_frame"):
        return ""
    parts = []
    for para in sh.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs)
        if not line:
            line = para.text
        parts.append(line)
    return "\n".join(parts).strip()


def get_prst_geo(sh):
    """Get the preset geometry name from the shape XML."""
    try:
        spPr = sh._element.find(".//" + qn("a:spPr"))
        if spPr is not None:
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                return prstGeom.get("prst")
    except Exception:
        pass
    return None


def get_line_endpoints(sh):
    """Extract begin/end coordinates for a connector or line shape."""
    try:
        el = sh._element
        xfrm = el.find(".//" + qn("a:xfrm"))
        if xfrm is None:
            return None, None

        flip_h = xfrm.get("flipH", "0") == "1"
        flip_v = xfrm.get("flipV", "0") == "1"

        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is None or ext is None:
            return None, None

        x = int(off.get("x", 0))
        y = int(off.get("y", 0))
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))

        x2 = x + cx if not flip_h else x - cx
        y2 = y + cy if not flip_v else y - cy

        return (emu_to_inches(x), emu_to_inches(y)), \
               (emu_to_inches(x2), emu_to_inches(y2))
    except Exception:
        return None, None


def get_arrowhead_info(sh):
    """Check if a line/connector has arrowhead markers and determine direction."""
    try:
        el = sh._element
        spPr = el.find(".//" + qn("a:spPr"))
        if spPr is None:
            return False, "—"

        ln = spPr.find(qn("a:ln"))
        if ln is None:
            return False, "—"

        head = ln.find(qn("a:headEnd"))
        tail = ln.find(qn("a:tailEnd"))

        has_head = head is not None and head.get("type", "none") != "none"
        has_tail = tail is not None and tail.get("type", "none") != "none"

        if has_tail and not has_head:
            return True, "→"
        elif has_head and not has_tail:
            return True, "←"
        elif has_head and has_tail:
            return True, "↔"
        else:
            return False, "—"
    except Exception:
        return False, "—"


def classify_shape_text(text, shape_type, prst_geo):
    """Classify a shape by its content and geometry into a semantic role."""
    if not text and shape_type not in ("LINE",):
        if prst_geo and prst_geo in PRST_GEO_TO_LABEL:
            return f"[{PRST_GEO_TO_LABEL[prst_geo]}]"
        return ""

    if text.strip().isdigit():
        return "[number]"
    if re.match(r"^\d{1,2}:\d{2}", text.strip()):
        return "[time]"
    if re.match(r"^\d+min$", text.strip()):
        return "[duration]"
    if re.match(r"^\d+%", text.strip()):
        return "[metric]"

    geo_label = PRST_GEO_TO_LABEL.get(prst_geo, "") if prst_geo else ""

    if shape_type == "TEXT_BOX":
        return "[text]"
    elif geo_label:
        return f"[{geo_label}]"
    elif shape_type == "AUTO_SHAPE":
        return "[box]"
    else:
        return f"[{shape_type}]"


def flatten_group(group_shape, group_path="", idx_counter=None):
    """Recursively flatten a group into its child shapes, tracking the path."""
    if idx_counter is None:
        idx_counter = [0]

    results = []
    current_path = group_path
    if group_shape.name:
        current_path = f"{group_path} > {group_shape.name}" if current_path else group_shape.name

    for child in group_shape.shapes:
        idx_counter[0] += 1
        if child.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(flatten_group(child, current_path, idx_counter))
        else:
            results.append((child, current_path, idx_counter[0]))
    return results


def collect_shapes(slide):
    """Collect all shapes from a slide, flattening groups."""
    shapes = []
    idx = 0
    for sh in slide.shapes:
        idx += 1
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child, gpath, child_idx in flatten_group(sh, "", None):
                shapes.append((child, gpath, child_idx))
        else:
            shapes.append((sh, "", idx))
    return shapes


def build_shape_info(sh, group_path, idx):
    """Build a ShapeInfo from a raw shape."""
    st = str(sh.shape_type).split("(")[0].strip() if sh.shape_type else "UNKNOWN"

    left = emu_to_inches(sh.left) if sh.left is not None else 0.0
    top = emu_to_inches(sh.top) if sh.top is not None else 0.0
    width = emu_to_inches(sh.width) if sh.width is not None else 0.0
    height = emu_to_inches(sh.height) if sh.height is not None else 0.0

    text = get_shape_text(sh)
    prst = get_prst_geo(sh)

    line_start = line_end = None
    has_arrow = False
    arrow_dir = "—"
    if st == "LINE" or (hasattr(sh, "shape_type") and
                        sh.shape_type == MSO_SHAPE_TYPE.LINE):
        line_start, line_end = get_line_endpoints(sh)
        has_arrow, arrow_dir = get_arrowhead_info(sh)

    is_table = sh.has_table
    table_rows = []
    if is_table:
        tbl = sh.table
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                cells.append(get_shape_text(cell).replace("\n", " "))
            table_rows.append(cells)

    picture_name = ""
    if "PICTURE" in st.upper():
        try:
            picture_name = sh.image.filename or "embedded"
        except Exception:
            picture_name = "embedded"

    return ShapeInfo(
        idx=idx, shape_type=st, prst_geo=prst, text=text,
        left=left, top=top, width=width, height=height,
        cx=left + width / 2, cy=top + height / 2,
        group_path=group_path,
        has_arrowhead=has_arrow, arrow_dir=arrow_dir,
        line_start=line_start, line_end=line_end,
        picture_name=picture_name,
        is_table=is_table, table_rows=table_rows,
        raw_shape=sh,
    )


# ── Spatial grid inference ───────────────────────────────────────────────

def cluster_centers(centers, threshold):
    if not centers:
        return []
    sorted_vals = sorted(centers)
    clusters = []
    current_cluster = [sorted_vals[0]]
    for i in range(1, len(sorted_vals)):
        if sorted_vals[i][0] - current_cluster[-1][0] <= threshold:
            current_cluster.append(sorted_vals[i])
        else:
            clusters.append(current_cluster)
            current_cluster = [sorted_vals[i]]
    clusters.append(current_cluster)
    return clusters


def assign_grid_positions(shapes, cfg):
    positioned = [s for s in shapes if s.shape_type not in ("LINE",) and s.width > 0]
    if not positioned:
        return {}, [], []

    y_centers = [(s.cy, i) for i, s in enumerate(positioned)]
    y_clusters = cluster_centers(y_centers, cfg["row_cluster_threshold"])

    row_map = {}
    row_labels = []
    for row_num, cluster in enumerate(y_clusters):
        avg_y = sum(v for v, _ in cluster) / len(cluster)
        row_labels.append(avg_y)
        for _, sidx in cluster:
            row_map[sidx] = row_num

    x_centers = [(s.cx, i) for i, s in enumerate(positioned)]
    x_clusters = cluster_centers(x_centers, cfg["col_cluster_threshold"])

    col_map = {}
    col_labels = []
    for col_num, cluster in enumerate(x_clusters):
        avg_x = sum(v for v, _ in cluster) / len(cluster)
        col_labels.append(avg_x)
        for _, sidx in cluster:
            col_map[sidx] = col_num

    return (row_map, col_map), row_labels, col_labels


# ── Connector inference ──────────────────────────────────────────────────

def find_nearest_shape(point, shapes, exclude_idx=None, proximity=0.35):
    if point is None:
        return None
    px, py = point

    candidates = []
    for i, s in enumerate(shapes):
        if exclude_idx is not None and i == exclude_idx:
            continue
        if s.shape_type in ("LINE",) or s.width <= 0:
            continue
        if (s.left - proximity <= px <= s.left + s.width + proximity and
                s.top - proximity <= py <= s.top + s.height + proximity):
            dx = max(s.left - px, 0, px - (s.left + s.width))
            dy = max(s.top - py, 0, py - (s.top + s.height))
            dist = math.sqrt(dx * dx + dy * dy)
            candidates.append((dist, i, s))

    if not candidates:
        for i, s in enumerate(shapes):
            if exclude_idx is not None and i == exclude_idx:
                continue
            if s.shape_type in ("LINE",) or s.width <= 0:
                continue
            dx = px - s.cx
            dy = py - s.cy
            dist = math.sqrt(dx * dx + dy * dy)
            candidates.append((dist, i, s))

    if not candidates:
        return None

    candidates.sort(key=lambda c: c[0])
    return candidates[0][2]


def infer_connectors(shapes, cfg):
    connections = []
    prox = cfg["connector_proximity"]
    for i, s in enumerate(shapes):
        if s.shape_type != "LINE" or s.line_start is None or s.line_end is None:
            continue

        source = find_nearest_shape(s.line_start, shapes, exclude_idx=i, proximity=prox)
        target = find_nearest_shape(s.line_end, shapes, exclude_idx=i, proximity=prox)

        sx, sy = s.line_start
        ex, ey = s.line_end
        dx = ex - sx
        dy = ey - sy

        if abs(dy) > abs(dx):
            dir_desc = "↓ down" if dy > 0 else "↑ up"
        elif abs(dx) > 0.1:
            dir_desc = "→ right" if dx > 0 else "← left"
        else:
            dir_desc = "·"

        connections.append({
            "source": source, "target": target,
            "dir": dir_desc, "arrow": s.arrow_dir,
            "has_arrow": s.has_arrowhead, "line_shape": s,
        })
    return connections


# ── Rendering ────────────────────────────────────────────────────────────

def truncate(text, max_len=60):
    if not text:
        return ""
    text = text.replace("\n", " | ")
    if len(text) > max_len:
        return text[:max_len - 1] + "…"
    return text


def render_table(table_rows):
    if not table_rows:
        return ""
    lines = []
    header = table_rows[0]
    lines.append("| " + " | ".join(c for c in header) + " |")
    lines.append("|" + "|".join("---" for _ in header) + "|")
    for row in table_rows[1:]:
        while len(row) < len(header):
            row.append("")
        lines.append("| " + " | ".join(c for c in row) + " |")
    return "\n".join(lines)


def render_spatial_grid(shapes, grid_data, row_labels, col_labels, cfg):
    (row_map, col_map) = grid_data
    if not row_labels or not col_labels:
        return ""

    n_rows = len(row_labels)
    n_cols = len(col_labels)

    cell_map = {}
    positioned_shapes = [s for s in shapes if s.shape_type not in ("LINE",) and s.width > 0]

    for i, s in enumerate(positioned_shapes):
        row = row_map.get(i)
        col = col_map.get(i)
        if row is None or col is None:
            continue

        role = classify_shape_text(s.text, s.shape_type, s.prst_geo)
        txt = truncate(s.text, cfg["truncate_text"])
        if role and txt:
            label = f"{role} {txt}"
        elif role:
            label = role
        elif txt:
            label = txt
        else:
            continue

        key = (row, col)
        if key not in cell_map:
            cell_map[key] = []
        cell_map[key].append(label)

    col_headers = [f"Col {c+1}\n(x≈{col_labels[c]:.1f}\")" for c in range(n_cols)]

    lines = []
    header = f"| Row\\\\Col | " + " | ".join(col_headers) + " |"
    lines.append(header)
    lines.append("|" + "---|" * (n_cols + 1))

    for r in range(n_rows):
        row_label = f"**R{r+1}**\n(y≈{row_labels[r]:.1f}\")"
        cells = []
        for c in range(n_cols):
            labels = cell_map.get((r, c), [])
            cells.append("<br>".join(labels) if labels else "")
        lines.append(f"| {row_label} | " + " | ".join(cells) + " |")

    return "\n".join(lines)


def render_connectors(connections):
    if not connections:
        return ""
    lines = ["**Connectors (inferred by spatial proximity):**"]
    for conn in connections:
        src_text = truncate(conn["source"].text, 30) if conn["source"] else "?"
        tgt_text = truncate(conn["target"].text, 30) if conn["target"] else "?"
        arrow = conn["arrow"] if conn["has_arrow"] else "—"
        dir_desc = conn["dir"]
        lines.append(f"- {arrow} `{src_text}` → `{tgt_text}` ({dir_desc})")
    return "\n".join(lines)


def render_picture(s):
    return f"[image: {s.picture_name or 'image'}, {s.width:.1f}\"×{s.height:.1f}\"]"


def find_slide_title(shapes_info):
    """Find slide title: first text box with text, near top."""
    text_shapes = [s for s in shapes_info
                   if s.shape_type == "TEXT_BOX" and s.text and s.top < 1.0]
    if text_shapes:
        text_shapes.sort(key=lambda s: s.top)
        return text_shapes[0].text.split("\n")[0]
    # Fallback: first shape with text
    for s in shapes_info:
        if s.text and s.top < 1.5:
            return s.text.split("\n")[0]
    return ""


def render_slide(slide, slide_num, shapes_info, cfg, slide_images=None, output_dir=None):
    lines = []
    lines.append("---")
    lines.append("")
    lines.append(f"## Slide {slide_num}")
    lines.append("")

    title = find_slide_title(shapes_info)

    n_total = len(shapes_info)
    n_tables = sum(1 for s in shapes_info if s.is_table)
    n_lines = sum(1 for s in shapes_info if s.shape_type == "LINE")
    n_pics = sum(1 for s in shapes_info if "PICTURE" in s.shape_type.upper())

    lines.append(f"**Title:** {title}" if title else "**Title:** (none)")
    lines.append(f"**Shapes:** {n_total} | {n_tables} tables | {n_lines} connectors | {n_pics} images")
    lines.append("")

    # 0. Embed slide image for complex slides (faithfully preserves tables/diagrams)
    threshold = cfg.get("image_embed_threshold", 0)
    if cfg.get("embed_slide_images", False) and slide_images and slide_num in slide_images:
        # Only embed image if slide has enough shapes to warrant it,
        # OR has a native table (complex layout likely)
        if n_total >= threshold or n_tables > 0:
            img_name = slide_images[slide_num]
            if output_dir:
                img_path = os.path.join(output_dir, img_name)
                data_uri = embed_image_base64(img_path, cfg.get("image_embed_max_kb", 500))
                if data_uri:
                    lines.append(f"![Slide {slide_num}]({data_uri})")
                else:
                    lines.append(f"![Slide {slide_num}]({img_name})")
            else:
                lines.append(f"![Slide {slide_num}]({img_name})")
            lines.append("")

    # 1. Native tables
    tables = [s for s in shapes_info if s.is_table]
    if tables:
        lines.append("### Tables")
        lines.append("")
        for ti, t in enumerate(tables, 1):
            lines.append(f"**Table {ti}:**")
            lines.append("")
            lines.append(render_table(t.table_rows))
            lines.append("")

    # 2. Text content — simple bullet list sorted by reading order (top→bottom, left→right)
    # This is the text-searchable supplement to the image. No spatial grid, no
    # connectors, no groups — those produce noise when the image is present.
    non_table = [s for s in shapes_info if not s.is_table and s.shape_type != "LINE"]
    text_shapes = [s for s in non_table if s.text and s.text.strip()]
    if text_shapes:
        lines.append("### Content")
        lines.append("")
        sorted_shapes = sorted(text_shapes, key=lambda s: (s.top, s.left))
        for s in sorted_shapes:
            txt = truncate(s.text, 200)
            if "PICTURE" in s.shape_type.upper():
                lines.append(f"- {render_picture(s)}")
            elif txt:
                lines.append(f"- {txt}")
        lines.append("")

    # 6. Speaker notes
    if cfg["include_notes"] and slide.has_notes_slide:
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, TypeError):
            notes_text = ""
        if notes_text:
            lines.append(f"**[Notes]** {notes_text}")
            lines.append("")

    return "\n".join(lines)


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Extract structured Markdown from a PPTX deck."
    )
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("-o", "--output", help="Output .md file path")
    parser.add_argument("--config", help="Path to config.yaml")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    cfg = load_config(args.config)

    # Determine output path
    if args.output:
        output_path = args.output
    else:
        input_stem = Path(args.input).stem
        output_dir = cfg.get("output_dir", "") or os.path.dirname(args.input)
        output_path = os.path.join(output_dir, input_stem + cfg["output_suffix"])

    prs = Presentation(args.input)

    output_lines = []
    # Title from core properties
    deck_title = ""
    try:
        deck_title = prs.core_properties.title or ""
    except Exception:
        pass
    if not deck_title:
        deck_title = Path(args.input).stem

    output_lines.append(f"# {deck_title} — Structured Extraction")
    output_lines.append("")
    output_lines.append(f"> Source: `{os.path.basename(args.input)}`")
    output_lines.append(f"> Slides: {len(prs.slides)}")
    output_lines.append(f"> Method: python-pptx structured extraction (spatial layout, connector topology, group nesting, shape semantics)")
    output_lines.append("")
    output_lines.append("---")
    output_lines.append("")

    total_shapes = 0
    total_tables = 0
    total_lines = 0
    total_pics = 0
    diagram_heavy = []

    # Export slide images if enabled
    slide_images = {}
    img_dir = None
    if cfg.get("embed_slide_images", False):
        # Determine image output directory (same as markdown output)
        if args.output:
            img_dir = str(Path(args.output).parent)
        else:
            img_dir = cfg.get("output_dir", "") or os.path.dirname(args.input)
        os.makedirs(img_dir, exist_ok=True)

        # First check for pre-existing slide images (e.g. exported separately)
        n_slides = len(prs.slides)
        for sn in range(1, n_slides + 1):
            for ext in (".png", ".jpg", ".jpeg"):
                candidate = os.path.join(img_dir, f"slide_{sn:02d}{ext}")
                if os.path.exists(candidate):
                    slide_images[sn] = f"slide_{sn:02d}{ext}"
                    break

        if len(slide_images) < n_slides:
            # Export missing images
            missing = [sn for sn in range(1, n_slides + 1) if sn not in slide_images]
            print(f"Exporting {len(missing)} missing slide images...", file=sys.stderr)
            new_images = export_slide_images(args.input, img_dir, slide_nums=missing)
            slide_images.update(new_images)

        if slide_images:
            print(f"  {len(slide_images)} slide images available in {img_dir}", file=sys.stderr)
        else:
            print("  (No slide images available — continuing with text-only extraction)", file=sys.stderr)

    for i, slide in enumerate(prs.slides, 1):
        raw_shapes = collect_shapes(slide)
        shapes_info = [build_shape_info(sh, gpath, idx) for sh, gpath, idx in raw_shapes]

        n = len(shapes_info)
        nt = sum(1 for s in shapes_info if s.is_table)
        nl = sum(1 for s in shapes_info if s.shape_type == "LINE")
        np_ = sum(1 for s in shapes_info if "PICTURE" in s.shape_type.upper())

        total_shapes += n
        total_tables += nt
        total_lines += nl
        total_pics += np_

        if n > cfg["diagram_heavy_threshold"]:
            diagram_heavy.append((i, n))

        slide_md = render_slide(slide, i, shapes_info, cfg, slide_images, img_dir)
        output_lines.append(slide_md)

    # Deck summary
    if cfg["include_deck_summary"]:
        summary = []
        summary.append("## Overview")
        summary.append("")
        summary.append("| Metric | Value |")
        summary.append("|---|---|")
        summary.append(f"| Total slides | {len(prs.slides)} |")
        summary.append(f"| Total shapes | {total_shapes} |")
        summary.append(f"| Native tables | {total_tables} |")
        summary.append(f"| Connectors | {total_lines} |")
        summary.append(f"| Images | {total_pics} |")
        summary.append(f"| Diagram-heavy slides (>{cfg['diagram_heavy_threshold']} shapes) | {len(diagram_heavy)}: {[d[0] for d in diagram_heavy]} |")
        summary.append("")
        summary.append("---")
        summary.append("")

        # Insert after first ---
        insert_pos = 0
        for i, line in enumerate(output_lines):
            if line.strip() == "---":
                insert_pos = i + 1
                break
        output_lines = output_lines[:insert_pos] + summary + output_lines[insert_pos:]

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(output_lines))

    print(f"Wrote {output_path}")
    print(f"  {len(prs.slides)} slides, {total_shapes} shapes, "
          f"{total_tables} tables, {total_lines} connectors, {total_pics} images")


if __name__ == "__main__":
    main()
